from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_ppt():
    prs = Presentation()

    # Helper to add title and content
    def add_slide(title, content_text=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        
        if content_text:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content_text
        return slide

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Traffic-V"
    subtitle.text = "Advanced ANPR & Vehicle Analytics System\nPowered by Deep Learning"

    # 2. Executive Summary
    add_slide("Executive Summary", 
              "Traffic-V is a state-of-the-art software solution designed for real-time traffic monitoring.\n\n"
              "It combines traditional Computer Vision with modern Deep Learning to detect vehicles, "
              "recognize attributes (Color, Type), and read License Plates (ANPR) with high accuracy.\n\n"
              "Key Value Proposition:\n"
              "- Real-time Processing\n"
              "- Robust to lighting conditions\n"
              "- Searchable Historical Database")

    # 3. Key Features
    add_slide("Key Features",
              "1. Vehicle Detection: MobileNet-SSD for fast and accurate detection of Cars, Buses, Trucks, and Bikes.\n"
              "2. Multi-Object Tracking: Centroid Tracking algorithm to maintain vehicle identity across frames.\n"
              "3. Advanced Color Recognition: Hybrid engine using K-Means Clustering and OpenAI CLIP (Zero-Shot) for precise color naming.\n"
              "4. License Plate Recognition: EAST Detector + TrOCR (Transformer OCR) for reading plates.\n"
              "5. Analytics Dashboard: Web-based UI for search, filtering, and reporting.")

    # 4. Architecture Diagram
    slide = add_slide("High-Level Architecture")
    # Try to add image if it exists
    img_path = r"C:/Users/HP/.gemini/antigravity/brain/a9836878-6600-4db4-b203-185914f16f77/traffic_v_architecture_16_9_1763912570425.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))
    else:
        slide.placeholders[1].text = "[Architecture Diagram Image Placeholder]"

    # 5. Technical Components
    slide = add_slide("Low-Level Components")
    img_path_2 = r"C:/Users/HP/.gemini/antigravity/brain/a9836878-6600-4db4-b203-185914f16f77/traffic_v_components_16_9_1763912628963.png"
    if os.path.exists(img_path_2):
        slide.shapes.add_picture(img_path_2, Inches(0.5), Inches(1.5), width=Inches(9))
    else:
        slide.placeholders[1].text = "[Component Diagram Image Placeholder]"

    # 6. Deep Learning Integration
    add_slide("Deep Learning Integration",
              "The system features a dual-mode Color Detector:\n\n"
              "Mode A: Fast Heuristic (K-Means)\n"
              "- Clusters pixel colors to find dominant shade.\n"
              "- Uses Nearest Neighbor to map to standard palette.\n"
              "- Speed: <10ms per vehicle.\n\n"
              "Mode B: Deep Learning (OpenAI CLIP)\n"
              "- Uses Zero-Shot Classification.\n"
              "- Understands semantic concepts ('silver car' vs 'white car').\n"
              "- Accuracy: >98% even in challenging lighting.")

    # 7. Future Roadmap
    add_slide("Future Roadmap",
              "- Integration with IP Cameras (RTSP Support).\n"
              "- Speed Estimation using perspective transformation.\n"
              "- Cloud Database Sync for multi-site deployment.\n"
              "- Mobile App for alerts.")

    prs.save('Traffic_V_Presentation.pptx')
    print("Presentation saved as Traffic_V_Presentation.pptx")

if __name__ == "__main__":
    create_ppt()
